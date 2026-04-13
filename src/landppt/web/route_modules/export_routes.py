"""
Export routes extracted from the legacy web router.
"""

from __future__ import annotations

import logging
import os
import tempfile
import time
import urllib.parse

from fastapi import APIRouter, Depends, HTTPException, Request
from fastapi.responses import FileResponse, JSONResponse

from ...auth.middleware import get_current_user_required
from ...database.models import User
from ...services.pdf_to_pptx_converter import get_pdf_to_pptx_converter
from ...services.pyppeteer_pdf_converter import get_pdf_converter
from ...utils.thread_pool import run_blocking_io
from .export_support import (
    ImagePPTXExportRequest,
    _generate_html_export_sync,
    _generate_pdf_with_pyppeteer,
    _is_standard_pptx_export_enabled,
    _prepare_html_for_file_based_export,
    _resolve_export_base_url,
)
from .support import ppt_service

router = APIRouter()

_TASK_PATH_KEYS = {"pdf_path", "pptx_path", "video_path", "audio_path"}


def _sanitize_task_mapping(payload: object) -> object:
    """Return a shallow-redacted view of task metadata/result for API responses."""
    if not isinstance(payload, dict):
        return payload
    return {
        key: value
        for key, value in payload.items()
        if key not in _TASK_PATH_KEYS
    }


def _get_task_owner_id(task) -> int | None:
    metadata = task.metadata if isinstance(task.metadata, dict) else {}
    owner_id = metadata.get("user_id")
    try:
        return int(owner_id) if owner_id is not None else None
    except (TypeError, ValueError):
        return None


def _ensure_task_access(task, user: User) -> None:
    if getattr(user, "is_admin", False):
        return

    owner_id = _get_task_owner_id(task)
    if owner_id is None or owner_id != int(user.id):
        raise HTTPException(status_code=404, detail="Task not found")


@router.get("/api/projects/{project_id}/export/pdf")
async def export_project_pdf(
    project_id: str,
    individual: bool = False,
    user: User = Depends(get_current_user_required)
):
    """Export project as PDF using Pyppeteer"""
    try:
        project = await ppt_service.project_manager.get_project(project_id, user_id=user.id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        # Check if we have slides data
        if not project.slides_data or len(project.slides_data) == 0:
            raise HTTPException(status_code=400, detail="PPT not generated yet")

        # Check if Pyppeteer is available
        pdf_converter = get_pdf_converter()
        if not pdf_converter.is_available():
            raise HTTPException(
                status_code=503,
                detail="PDF generation service unavailable. Please ensure Pyppeteer is installed: pip install pyppeteer"
            )

        # Create temp file in thread pool to avoid blocking
        temp_pdf_path = await run_blocking_io(
            lambda: tempfile.NamedTemporaryFile(suffix='.pdf', delete=False).name
        )

        logging.info("Generating PDF with Pyppeteer")
        success = await _generate_pdf_with_pyppeteer(project, temp_pdf_path, individual)

        if not success:
            # Clean up temp file and raise error
            await run_blocking_io(lambda: os.unlink(temp_pdf_path) if os.path.exists(temp_pdf_path) else None)
            raise HTTPException(status_code=500, detail="PDF generation failed")

        # Return PDF file
        logging.info("PDF generated successfully using Pyppeteer")
        safe_filename = urllib.parse.quote(f"{project.topic}_PPT.pdf", safe='')

        # 使用BackgroundTask来清理临时文件
        from starlette.background import BackgroundTask

        def cleanup_temp_file():
            try:
                os.unlink(temp_pdf_path)
            except:
                pass

        return FileResponse(
            temp_pdf_path,
            media_type="application/pdf",
            headers={
                "Content-Disposition": f"attachment; filename*=UTF-8''{safe_filename}",
                "X-PDF-Generator": "Pyppeteer"
            },
            background=BackgroundTask(cleanup_temp_file)
        )

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/api/projects/{project_id}/export/pdf/individual")
async def export_project_pdf_individual(
    project_id: str,
    user: User = Depends(get_current_user_required)
):
    """Export project as individual PDF files for each slide"""
    return await export_project_pdf(project_id, individual=True, user=user)


@router.post("/api/projects/{project_id}/export/pdf/async")
async def export_project_pdf_async(
    project_id: str,
    user: User = Depends(get_current_user_required)
):
    """Export project as PDF using background task with progress tracking.
    
    Returns immediately with task_id. Poll /api/landppt/tasks/{task_id} for progress.
    Download result from /api/landppt/tasks/{task_id}/download when completed.
    """
    try:
        project = await ppt_service.project_manager.get_project(project_id, user_id=user.id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        # Check if we have slides data
        if not project.slides_data or len(project.slides_data) == 0:
            raise HTTPException(status_code=400, detail="PPT not generated yet")

        # Check if Pyppeteer/Playwright is available
        pdf_converter = get_pdf_converter()
        if not pdf_converter.is_available():
            raise HTTPException(
                status_code=503,
                detail="PDF generation service unavailable. Please ensure Playwright is installed."
            )

        from ...services.background_tasks import get_task_manager, TaskStatus

        task_manager = get_task_manager()
        
        # Check if there's already an active PDF export task for this project
        existing_task = await task_manager.find_active_task_async(
            task_type="pdf_generation",
            metadata_filter={"project_id": project_id, "user_id": user.id}
        )
        if existing_task:
            logging.info(f"Project {project_id} already has an active PDF export task: {existing_task.task_id}")
            return JSONResponse(
                status_code=409,
                content={
                    "status": "already_processing",
                    "task_id": existing_task.task_id,
                    "message": "PDF generation is already in progress for this project",
                    "polling_endpoint": f"/api/landppt/tasks/{existing_task.task_id}"
                }
            )

        # Create temp file for output
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf_file:
            temp_pdf_path = temp_pdf_file.name

        # Define PDF generation task
        async def pdf_generation_task():
            """PDF generation background task with progress updates."""
            try:
                total_slides = len(project.slides_data)
                
                # Progress callback for the PDF converter
                async def progress_callback(current_index: int, total: int):
                    progress = (current_index / total) * 100 if total > 0 else 0
                    task_manager.update_task_status(
                        task_id,
                        TaskStatus.RUNNING,
                        progress=progress
                    )
                
                logging.info(f"Starting background PDF generation for {total_slides} slides")
                
                success = await _generate_pdf_with_pyppeteer(project, temp_pdf_path, individual=False)
                
                if success and os.path.exists(temp_pdf_path):
                    return {
                        "success": True,
                        "pdf_path": temp_pdf_path,
                        "project_topic": project.topic,
                        "slide_count": total_slides
                    }
                else:
                    # Clean up on failure
                    try:
                        if os.path.exists(temp_pdf_path):
                            os.unlink(temp_pdf_path)
                    except:
                        pass
                    return {
                        "success": False,
                        "error": "PDF generation failed"
                    }
                    
            except Exception as e:
                logging.error(f"PDF generation task error: {e}")
                import traceback
                traceback.print_exc()
                # Clean up on error
                try:
                    if os.path.exists(temp_pdf_path):
                        os.unlink(temp_pdf_path)
                except:
                    pass
                raise

        # Submit task to background
        task_id = task_manager.submit_task(
            task_type="pdf_generation",
            func=pdf_generation_task,
            metadata={
                "project_id": project_id,
                "project_topic": project.topic,
                "slide_count": len(project.slides_data),
                "user_id": user.id,
            }
        )

        logging.info(f"PDF generation task started: {task_id}")

        return JSONResponse({
            "status": "processing",
            "task_id": task_id,
            "message": "PDF generation started in background",
            "polling_endpoint": f"/api/landppt/tasks/{task_id}",
            "slide_count": len(project.slides_data)
        })

    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Failed to start PDF generation task: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/api/projects/{project_id}/export/pptx")
async def export_project_pptx(
    project_id: str,
    user: User = Depends(get_current_user_required)
):
    """Export project as PPTX by first generating PDF then converting to PowerPoint.
    
    Returns immediately with task_id. Poll /api/landppt/tasks/{task_id} for progress.
    Download result from /api/landppt/tasks/{task_id}/download when completed.
    """
    try:
        project = await ppt_service.project_manager.get_project(project_id, user_id=user.id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        if not await _is_standard_pptx_export_enabled():
            raise HTTPException(
                status_code=403,
                detail="Apryse PPTX export is disabled by administrator or not configured."
            )

        # Check if we have slides data
        if not project.slides_data or len(project.slides_data) == 0:
            raise HTTPException(status_code=400, detail="PPT not generated yet")

        # Pre-load user's license key (async, fast operation)
        converter = get_pdf_to_pptx_converter()
        await converter.set_user_id_async(user.id)
        
        # NOTE: Skip blocking is_available() check here - it will be done in background task
        # This prevents SDK download from blocking the request

        from ...services.background_tasks import get_task_manager, TaskStatus

        task_manager = get_task_manager()
        
        # 检查是否已有该项目的活跃导出任务（防止重复请求）
        existing_task = await task_manager.find_active_task_async(
            task_type="pdf_to_pptx_conversion",
            metadata_filter={"project_id": project_id, "user_id": user.id}
        )
        if existing_task:
            logging.info(f"Project {project_id} already has an active PPTX export task: {existing_task.task_id}")
            return JSONResponse(
                status_code=409,
                content={
                    "status": "already_processing",
                    "task_id": existing_task.task_id,
                    "message": "PPTX conversion is already in progress for this project",
                    "polling_endpoint": f"/api/landppt/tasks/{existing_task.task_id}"
                }
            )

        # 创建临时文件路径
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf_file:
            temp_pdf_path = temp_pdf_file.name
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_pptx_file:
            temp_pptx_path = temp_pptx_file.name

        # 定义完整的转换任务函数（PDF生成 + PPTX转换）
        async def pdf_to_pptx_task():
            """Complete PPTX export task: PDF generation + conversion (runs in background)."""
            try:
                total_slides = len(project.slides_data) if project.slides_data else 0
                
                # Step 0: Check SDK availability (this may download SDK, blocking but in background)
                logging.info("Step 0: Checking SDK and PDF converter availability")
                task_manager.update_task_status(
                    task_id,
                    TaskStatus.RUNNING,
                    progress=1
                )
                
                # Check PDF converter availability
                pdf_converter = get_pdf_converter()
                if not pdf_converter.is_available():
                    return {
                        "success": False,
                        "error": "PDF generation service unavailable. Please ensure Playwright is installed."
                    }
                
                # Check PPTX converter availability (may trigger SDK download)
                if not converter.is_available():
                    return {
                        "success": False,
                        "error": "PPTX conversion service unavailable. Please ensure Apryse SDK is installed and licensed."
                    }
                
                task_manager.update_task_status(
                    task_id,
                    TaskStatus.RUNNING,
                    progress=5
                )
                
                # Step 1: 生成 PDF
                logging.info(f"Step 1: Generating PDF for {total_slides} slides")
                pdf_success = await _generate_pdf_with_pyppeteer(project, temp_pdf_path, individual=False)
                
                if not pdf_success:
                    # 清理临时文件
                    try:
                        if os.path.exists(temp_pdf_path):
                            os.unlink(temp_pdf_path)
                    except:
                        pass
                    return {
                        "success": False,
                        "error": "PDF generation failed"
                    }
                
                logging.info("Step 1 completed: PDF generated successfully")
                task_manager.update_task_status(
                    task_id,
                    TaskStatus.RUNNING,
                    progress=50
                )
                
                # Step 2: 转换 PDF 到 PPTX
                logging.info("Step 2: Converting PDF to PPTX")
                success, result = await converter.convert_pdf_to_pptx_async(
                    temp_pdf_path,
                    temp_pptx_path
                )
                if success:
                    task_manager.update_task_status(
                        task_id,
                        TaskStatus.RUNNING,
                        progress=90
                    )
                    
                    # 转换成功后，添加演讲稿到备注
                    try:
                        from pptx import Presentation
                        from ...services.speech_script_repository import SpeechScriptRepository

                        # 获取演讲稿数据
                        repo = SpeechScriptRepository()
                        scripts_list = await repo.get_current_speech_scripts_by_project(project_id)
                        speech_scripts = {script.slide_index: script.script_content for script in scripts_list}
                        repo.close()

                        if len(speech_scripts) > 0:
                            # 打开生成的PPTX文件
                            prs = Presentation(temp_pptx_path)

                            # 为每张幻灯片添加演讲稿备注
                            for i, slide in enumerate(prs.slides):
                                if i in speech_scripts:
                                    notes_slide = slide.notes_slide
                                    text_frame = notes_slide.notes_text_frame
                                    text_frame.text = speech_scripts[i]
                                    logging.info(f"Added speech script to slide {i+1} notes")

                            # 保存修改后的PPTX
                            prs.save(temp_pptx_path)
                            logging.info(f"Added {len(speech_scripts)} speech scripts to PPTX notes")
                    except Exception as e:
                        logging.warning(f"Failed to add speech scripts to PPTX: {e}")
                        # 继续执行，即使添加演讲稿失败也返回PPTX

                    return {
                        "success": True,
                        "pptx_path": temp_pptx_path,
                        "pdf_path": temp_pdf_path
                    }
                else:
                    # 清理临时文件
                    try:
                        if os.path.exists(temp_pdf_path):
                            os.unlink(temp_pdf_path)
                        if os.path.exists(temp_pptx_path):
                            os.unlink(temp_pptx_path)
                    except:
                        pass
                    return {
                        "success": False,
                        "error": result
                    }
            except Exception as e:
                logging.error(f"PPTX export task error: {e}")
                import traceback
                traceback.print_exc()
                # 清理临时文件
                try:
                    if os.path.exists(temp_pdf_path):
                        os.unlink(temp_pdf_path)
                    if os.path.exists(temp_pptx_path):
                        os.unlink(temp_pptx_path)
                except:
                    pass
                return {
                    "success": False,
                    "error": str(e)
                }


        # 提交后台任务
        task_id = task_manager.submit_task(
            task_type="pdf_to_pptx_conversion",
            func=pdf_to_pptx_task,
            metadata={
                "project_id": project_id,
                "project_topic": project.topic,
                "slide_count": len(project.slides_data),
                "pdf_path": temp_pdf_path,
                "pptx_path": temp_pptx_path,
                "user_id": user.id,
            }
        )

        logging.info(f"PPTX export task started: {task_id}")

        # 立即返回任务ID，不等待任务完成
        return JSONResponse({
            "status": "processing",
            "task_id": task_id,
            "message": "PPTX export started in background (PDF generation + conversion)",
            "polling_endpoint": f"/api/landppt/tasks/{task_id}",
            "slide_count": len(project.slides_data)
        })

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/api/projects/{project_id}/export/pptx-images")
async def export_project_pptx_from_images(
    project_id: str,
    payload: ImagePPTXExportRequest,
    http_request: Request,
    user: User = Depends(get_current_user_required)
):
    """Export project as PPTX using high-quality Playwright screenshots"""
    try:
        from io import BytesIO
        from pptx import Presentation
        from pptx.util import Inches

        project = await ppt_service.project_manager.get_project(project_id, user_id=user.id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        # 验证是否有幻灯片数据
        slides = getattr(payload, 'slides', None)
        if not slides or len(slides) == 0:
            raise HTTPException(status_code=400, detail="No slides provided")

        export_base_url = _resolve_export_base_url(http_request)

        # 检查Playwright是否可用
        pdf_converter = get_pdf_converter()
        if not pdf_converter.is_available():
            raise HTTPException(
                status_code=503,
                detail="Screenshot service unavailable. Please ensure Playwright is installed."
            )

        # 创建后台任务
        from ...services.background_tasks import get_task_manager, TaskStatus
        task_manager = get_task_manager()

        # 创建临时目录和PPTX文件路径
        temp_dir = tempfile.mkdtemp()
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_pptx_file:
            temp_pptx_path = temp_pptx_file.name

        # 定义HTML到图片到PPTX的任务函数
        async def html_to_pptx_task():
            """使用Playwright截图并生成PPTX"""
            screenshot_paths = []
            try:
                logging.info(f"Starting screenshot-based PPTX export for {len(slides)} slides")
                total_slides = len(slides)

                def update_export_progress(progress: float, message: str):
                    try:
                        task = task_manager.tasks.get(task_id)
                        if task is not None:
                            task.metadata["progress_message"] = message
                        task_manager.update_task_status(
                            task_id,
                            TaskStatus.RUNNING,
                            progress=max(0.0, min(99.0, float(progress)))
                        )
                    except Exception as progress_error:
                        logging.debug(f"Failed to update screenshot export progress: {progress_error}")

                # 第1步：获取演讲稿数据
                update_export_progress(5, "正在加载图片导出任务...")
                speech_scripts = {}
                try:
                    from ...services.speech_script_repository import SpeechScriptRepository
                    repo = SpeechScriptRepository()
                    scripts_list = await repo.get_current_speech_scripts_by_project(project_id)
                    # 构建幻灯片索引到演讲稿的映射
                    for script in scripts_list:
                        speech_scripts[script.slide_index] = script.script_content
                    repo.close()
                    logging.info(f"Loaded {len(speech_scripts)} speech scripts for slides")
                except Exception as e:
                    logging.warning(f"Failed to load speech scripts: {e}")
                    # 继续执行，即使没有演讲稿也可以生成PPTX

                # 第2步：为每张幻灯片创建临时HTML文件
                update_export_progress(12, "正在准备幻灯片HTML文件...")
                html_files = []
                for i, slide in enumerate(slides):
                    html_file = os.path.join(temp_dir, f"slide_{i}.html")
                    prepared_html = _prepare_html_for_file_based_export(
                        slide.get('html_content', ''),
                        export_base_url
                    )
                    with open(html_file, 'w', encoding='utf-8') as f:
                        f.write(prepared_html)
                    html_files.append(html_file)
                    prep_progress = 12 + ((i + 1) / max(total_slides, 1)) * 13
                    update_export_progress(prep_progress, f"正在准备第 {i + 1}/{total_slides} 张幻灯片...")

                # 第3步：使用Playwright对每张幻灯片进行截图
                update_export_progress(25, "正在渲染幻灯片图片...")
                for i, html_file in enumerate(html_files):
                    screenshot_path = os.path.join(temp_dir, f"slide_{i}.png")

                    # 使用PDF converter的截图功能
                    success = await pdf_converter.screenshot_html(
                        html_file,
                        screenshot_path,
                        width=1280,
                        height=720,
                        optimize_for_static=True,
                        stability_checks=1,
                        stability_interval=0.2,
                    )

                    if success:
                        screenshot_paths.append(screenshot_path)
                        logging.info(f"Screenshot {i+1}/{len(html_files)} completed")
                    else:
                        logging.warning(f"Screenshot {i+1} failed, skipping")

                    screenshot_progress = 25 + ((i + 1) / max(len(html_files), 1)) * 55
                    update_export_progress(
                        screenshot_progress,
                        f"正在渲染第 {i + 1}/{len(html_files)} 张幻灯片图片..."
                    )

                if len(screenshot_paths) == 0:
                    raise Exception("No screenshots were generated")

                # 第4步：将截图转换为PPTX
                logging.info("Creating PPTX from screenshots...")
                update_export_progress(82, "正在组装PPTX文件...")
                prs = Presentation()

                # 设置幻灯片尺寸为16:9
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(5.625)

                for i, screenshot_path in enumerate(screenshot_paths):
                    # 添加空白幻灯片
                    blank_slide_layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(blank_slide_layout)

                    # 添加截图，填充整个幻灯片
                    left = Inches(0)
                    top = Inches(0)
                    width = prs.slide_width
                    height = prs.slide_height

                    slide.shapes.add_picture(screenshot_path, left, top, width=width, height=height)

                    # 如果该幻灯片有演讲稿，添加到备注中
                    if i in speech_scripts:
                        notes_slide = slide.notes_slide
                        text_frame = notes_slide.notes_text_frame
                        text_frame.text = speech_scripts[i]
                        logging.info(f"Added speech script to slide {i+1} notes")

                    pptx_progress = 82 + ((i + 1) / max(len(screenshot_paths), 1)) * 13
                    update_export_progress(
                        pptx_progress,
                        f"正在写入第 {i + 1}/{len(screenshot_paths)} 页到PPTX..."
                    )

                # 保存PPTX文件
                update_export_progress(97, "正在保存PPTX文件...")
                prs.save(temp_pptx_path)
                logging.info(f"PPTX saved to {temp_pptx_path}")

                return {
                    "success": True,
                    "pptx_path": temp_pptx_path
                }

            except Exception as e:
                logging.error(f"HTML to PPTX conversion failed: {e}")
                import traceback
                traceback.print_exc()
                return {
                    "success": False,
                    "error": str(e)
                }
            finally:
                # 清理临时HTML和截图文件
                try:
                    import shutil
                    if os.path.exists(temp_dir):
                        shutil.rmtree(temp_dir)
                        logging.info(f"Cleaned up temp directory: {temp_dir}")
                except Exception as cleanup_error:
                    logging.warning(f"Failed to cleanup temp directory: {cleanup_error}")

        # 提交后台任务
        task_id = task_manager.submit_task(
            task_type="html_to_pptx_screenshot",
            func=html_to_pptx_task,
            metadata={
                "project_id": project_id,
                "project_topic": project.topic,
                "slide_count": len(slides),
                "pptx_path": temp_pptx_path,
                "progress_message": "图片导出任务已创建，等待后台执行...",
                "user_id": user.id,
            }
        )

        # 立即返回任务ID
        return JSONResponse({
            "status": "processing",
            "task_id": task_id,
            "message": "PPTX generation with screenshots started in background",
            "polling_endpoint": f"/api/landppt/tasks/{task_id}"
        })

    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"PPTX screenshot export error: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/api/landppt/tasks/{task_id}")
async def get_task_status(
    task_id: str,
    user: User = Depends(get_current_user_required),
):
    """查询后台任务状态"""
    from ...services.background_tasks import get_task_manager

    task_manager = get_task_manager()
    # Use async version to check Valkey cache (for cross-worker task lookup)
    task = await task_manager.get_task_async(task_id)

    if not task:
        raise HTTPException(status_code=404, detail="Task not found")
    _ensure_task_access(task, user)


    response = {
        "task_id": task.task_id,
        "task_type": task.task_type,
        "status": task.status.value,
        "progress": task.progress,
        "created_at": task.created_at.isoformat(),
        "updated_at": task.updated_at.isoformat(),
        "metadata": _sanitize_task_mapping(task.metadata)
    }

    progress_message = task.metadata.get("progress_message") if isinstance(task.metadata, dict) else None
    if progress_message:
        response["message"] = progress_message

    # 如果任务完成，添加结果信息
    if task.status.value == "completed" and task.result:
        response["result"] = _sanitize_task_mapping(task.result)
        # 如果是可下载的导出任务，提供下载链接
        if isinstance(task.result, dict):
            if task.task_type in (
                "pdf_to_pptx_conversion",
                "pdf_generation",
                "html_to_pptx_screenshot",
                "narration_audio_export",
                "narration_video_export",
            ):
                # Provide the download endpoint for export tasks. The download handler will validate existence.
                response["download_url"] = f"/api/landppt/tasks/{task_id}/download"
            # Surface error details even if status is "completed" (legacy behavior).
            if task.result.get("success") is False:
                response["error"] = task.result.get("error") or task.result.get("message") or task.error

    # 如果任务失败，添加错误信息
    if task.status.value == "failed":
        response["error"] = task.error

    return JSONResponse(response)


@router.get("/api/landppt/tasks/{task_id}/download")
async def download_task_result(
    task_id: str,
    user: User = Depends(get_current_user_required),
):
    """下载任务结果文件（支持PDF、PPTX、讲解音频ZIP、讲解视频MP4）"""
    from ...services.background_tasks import get_task_manager, TaskStatus
    from starlette.background import BackgroundTask

    task_manager = get_task_manager()
    # Use async version to check Valkey cache (for cross-worker task lookup)
    task = await task_manager.get_task_async(task_id)

    if not task:
        raise HTTPException(status_code=404, detail="Task not found")
    _ensure_task_access(task, user)


    if task.status != TaskStatus.COMPLETED:
        raise HTTPException(status_code=400, detail=f"Task not completed yet (status: {task.status.value})")

    result = task.result if isinstance(task.result, dict) else {}
    pptx_path = result.get("pptx_path")
    pdf_path = result.get("pdf_path")
    video_path = result.get("video_path")
    audio_path = result.get("audio_path")
    project_topic = task.metadata.get("project_topic", "PPT")

    # 清理临时文件的后台任务
    def cleanup_temp_files():
        try:
            if pdf_path and os.path.exists(pdf_path):
                os.unlink(pdf_path)
        except:
            pass
        try:
            if pptx_path and os.path.exists(pptx_path):
                os.unlink(pptx_path)
        except:
            pass
        try:
            if audio_path and os.path.exists(audio_path):
                os.unlink(audio_path)
        except:
            pass

    if task.task_type == "narration_audio_export":
        if audio_path and os.path.exists(audio_path):
            ext = (os.path.splitext(audio_path)[1] or ".zip").lower()
            media_type = {
                ".zip": "application/zip",
                ".mp3": "audio/mpeg",
                ".wav": "audio/wav",
                ".flac": "audio/flac",
                ".m4a": "audio/mp4",
                ".ogg": "audio/ogg",
            }.get(ext, "application/octet-stream")
            safe_filename = urllib.parse.quote(
                f"{project_topic}_讲解音频_{task.metadata.get('language', 'zh')}{ext}",
                safe="",
            )
            return FileResponse(
                audio_path,
                media_type=media_type,
                headers={
                    "Content-Disposition": f"attachment; filename*=UTF-8''{safe_filename}",
                    "X-Export-Method": "Narration-Audio",
                },
                background=BackgroundTask(cleanup_temp_files),
            )

        if isinstance(result, dict) and result.get("error"):
            raise HTTPException(status_code=400, detail=str(result.get("error")))
        if task.error:
            raise HTTPException(status_code=400, detail=str(task.error))
        raise HTTPException(status_code=404, detail="Result file not found")

    # Special-case narration video: allow download even if legacy result didn't persist cleanly.
    if task.task_type == "narration_video_export":
        if not video_path:
            video_path = task.metadata.get("video_path")

        if not video_path:
            try:
                project_id = task.metadata.get("project_id")
                language = task.metadata.get("language", "zh")
                fps = task.metadata.get("fps")
                out_dir = os.path.join("uploads", "narration_videos", str(project_id), str(language))
                if os.path.isdir(out_dir):
                    candidates = []
                    for name in os.listdir(out_dir):
                        if not name.lower().endswith(".mp4"):
                            continue
                        if fps and f"{int(fps)}fps" not in name:
                            continue
                        candidates.append(os.path.join(out_dir, name))
                    if candidates:
                        candidates.sort(key=lambda p: os.path.getmtime(p), reverse=True)
                        video_path = candidates[0]
            except Exception:
                pass

        if video_path and os.path.exists(video_path):
            safe_filename = urllib.parse.quote(
                f"{project_topic}_narration_{task.metadata.get('language','zh')}.mp4",
                safe="",
            )
            # Videos are stored under uploads/; keep them by default (no temp cleanup).
            return FileResponse(
                video_path,
                media_type="video/mp4",
                headers={
                    "Content-Disposition": f"attachment; filename*=UTF-8''{safe_filename}",
                    "X-Export-Method": "Narration-Video",
                },
            )

        # Provide a more helpful message when the result isn't downloadable.
        if isinstance(result, dict) and result.get("error"):
            raise HTTPException(status_code=400, detail=str(result.get("error")))
        if task.error:
            raise HTTPException(status_code=400, detail=str(task.error))
        raise HTTPException(status_code=404, detail="Result file not found")

    if not result or not result.get("success"):
        raise HTTPException(status_code=400, detail="Task failed or no result available")

    # 根据任务类型返回不同的文件
    if task.task_type == "pdf_generation" and pdf_path and os.path.exists(pdf_path):
        # PDF生成任务：返回PDF文件
        safe_filename = urllib.parse.quote(f"{project_topic}_PPT.pdf", safe='')
        return FileResponse(
            pdf_path,
            media_type="application/pdf",
            headers={
                "Content-Disposition": f"attachment; filename*=UTF-8''{safe_filename}",
                "X-Conversion-Method": "PDF-Background"
            },
            background=BackgroundTask(cleanup_temp_files)
        )
    elif pptx_path and os.path.exists(pptx_path):
        # PPTX转换任务：返回PPTX文件
        safe_filename = urllib.parse.quote(f"{project_topic}_PPT.pptx", safe='')
        return FileResponse(
            pptx_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": f"attachment; filename*=UTF-8''{safe_filename}",
                "X-Conversion-Method": "PDF-to-PPTX-Background"
            },
            background=BackgroundTask(cleanup_temp_files)
        )
    else:
        raise HTTPException(status_code=404, detail="Result file not found")


@router.get("/api/projects/{project_id}/export/html")
async def export_project_html(
    project_id: str,
    user: User = Depends(get_current_user_required)
):
    """Export project as HTML ZIP package with slideshow index"""
    try:
        project = await ppt_service.project_manager.get_project(project_id, user_id=user.id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        # Check if we have slides data
        if not project.slides_data or len(project.slides_data) == 0:
            raise HTTPException(status_code=400, detail="PPT not generated yet")

        # Create temporary directory and generate files in thread pool
        zip_content = await run_blocking_io(_generate_html_export_sync, project)

        # URL encode the filename to handle Chinese characters
        zip_filename = f"{project.topic}_PPT.zip"
        safe_filename = urllib.parse.quote(zip_filename, safe='')

        from fastapi.responses import Response
        return Response(
            content=zip_content,
            media_type="application/zip",
            headers={
                "Content-Disposition": f"attachment; filename*=UTF-8''{safe_filename}"
            }
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
